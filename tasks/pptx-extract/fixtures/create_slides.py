#!/usr/bin/env python3
"""Create a sample PPTX for extraction task."""
import sys
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Annual Sales Review"
slide.placeholders[1].text = "Fiscal Year 2025"

# Slide 2: Revenue by Region
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Revenue by Region"
tf = slide.placeholders[1].text_frame
tf.text = "North America: $4.2M"
tf.add_paragraph().text = "Europe: $2.8M"
tf.add_paragraph().text = "Asia Pacific: $1.9M"
tf.add_paragraph().text = "Latin America: $0.6M"

# Slide 3: Key Metrics
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Metrics"
tf = slide.placeholders[1].text_frame
tf.text = "Total Revenue: $9.5M"
tf.add_paragraph().text = "YoY Growth: 23%"
tf.add_paragraph().text = "Customer Count: 1,247"
tf.add_paragraph().text = "Churn Rate: 4.2%"
tf.add_paragraph().text = "NPS Score: 72"

# Slide 4: Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
tf = slide.placeholders[1].text_frame
tf.text = "Expand APAC sales team by Q2"
tf.add_paragraph().text = "Launch enterprise tier in Europe"
tf.add_paragraph().text = "Target 30% YoY growth for FY2026"

prs.save(sys.argv[1] if len(sys.argv) > 1 else "review.pptx")
