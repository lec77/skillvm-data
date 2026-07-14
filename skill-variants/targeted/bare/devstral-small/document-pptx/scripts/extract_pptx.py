#!/usr/bin/env python3
"""Extract structured data from a PPTX file to JSON.

Usage:
    python3 scripts/extract_pptx.py input.pptx output.json

Extracts: title, subtitle, slide_count, slides (title + bullets),
and parses financial data (total_revenue, regions, yoy_growth).
"""
import json, sys, re
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

input_path = sys.argv[1]
output_path = sys.argv[2] if len(sys.argv) > 2 else "slides_data.json"

prs = Presentation(input_path)
slides = []
title = ""
subtitle = ""
regions = []
total_revenue = None
yoy_growth = None

for i, slide in enumerate(prs.slides):
    slide_title = slide.shapes.title.text.strip() if slide.shapes.title else ""
    bullets = []

    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            if i == 0:
                subtitle = shape.text_frame.text.strip()
            else:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        bullets.append(t)

    if i == 0:
        title = slide_title

    slides.append({"title": slide_title, "bullets": bullets})

    # Parse financial data from bullets
    for b in bullets:
        # Region: "North America: $4.2M"
        m = re.match(r'^(.+?):\s*\$?([\d.]+)\s*M', b)
        if m and 'revenue' not in m.group(1).lower() and 'growth' not in m.group(1).lower():
            regions.append({"name": m.group(1).strip(), "revenue": float(m.group(2))})

        # Total Revenue: "$9.5M"
        m2 = re.match(r'Total Revenue:\s*\$?([\d.]+)\s*M', b)
        if m2:
            total_revenue = float(m2.group(1))

        # YoY Growth: "23%"
        m3 = re.match(r'YoY Growth:\s*([\d.]+)\s*%', b)
        if m3:
            yoy_growth = int(float(m3.group(1)))

result = {
    "title": title,
    "subtitle": subtitle,
    "slide_count": len(prs.slides),
    "slides": slides,
}

if total_revenue is not None:
    result["total_revenue"] = total_revenue
if regions:
    result["regions"] = regions
if yoy_growth is not None:
    result["yoy_growth"] = yoy_growth

with open(output_path, "w") as f:
    json.dump(result, f, indent=2)

print(f"Extracted {len(slides)} slides to {output_path}")
